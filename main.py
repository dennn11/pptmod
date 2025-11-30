#!/usr/bin/env python3
"""
PowerPoint Text Modifier CLI
Modifies text in PowerPoint presentations (.ppt and .pptx) based on a config file.
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
import win32com.client
import os


def load_config(config_path: str) -> Dict[str, str]:
    """Load text replacements from a JSON config file."""
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
        
        if 'replacements' not in config:
            print("Error: Config file must contain 'replacements' key", file=sys.stderr)
            sys.exit(1)
        
        return config['replacements']
    except FileNotFoundError:
        print(f"Error: Config file '{config_path}' not found", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in config file: {e}", file=sys.stderr)
        sys.exit(1)


def modify_pptx(input_file: str, output_file: str, replacements: Dict[str, str]) -> None:
    """Modify a .pptx file using python-pptx library."""
    try:
        prs = Presentation(input_file)
        replacement_count = 0
        
        # Iterate through all slides
        for slide in prs.slides:
            # Check all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Iterate through paragraphs and runs
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old_text, new_text in replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    replacement_count += 1
                
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    for old_text, new_text in replacements.items():
                                        if old_text in run.text:
                                            run.text = run.text.replace(old_text, new_text)
                                            replacement_count += 1
        
        prs.save(output_file)
        print(f"Successfully modified {input_file} -> {output_file}")
        print(f"Made {replacement_count} text replacements")
    
    except Exception as e:
        print(f"Error modifying .pptx file: {e}", file=sys.stderr)
        sys.exit(1)


def export_to_pdf(input_file: str, output_pdf: str) -> None:
    """Export a PowerPoint file to PDF using COM automation (Windows only)."""
    try:
        # Convert to absolute paths
        input_path = os.path.abspath(input_file)
        output_path = os.path.abspath(output_pdf)
        
        # Initialize PowerPoint
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        
        # Open the presentation
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        
        # Export to PDF (32 = ppSaveAsPDF)
        presentation.SaveAs(output_path, 32)
        
        # Close and cleanup
        presentation.Close()
        powerpoint.Quit()
        
        print(f"Successfully exported to PDF: {output_pdf}")
    
    except Exception as e:
        print(f"Error exporting to PDF: {e}", file=sys.stderr)
        print("Note: PDF export requires PowerPoint to be installed on Windows", file=sys.stderr)
        raise


def modify_ppt(input_file: str, output_file: str, replacements: Dict[str, str]) -> None:
    """Modify a .ppt file using COM automation (Windows only)."""
    try:
        # Convert to absolute paths
        input_path = os.path.abspath(input_file)
        output_path = os.path.abspath(output_file)
        
        # Initialize PowerPoint
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        
        # Open the presentation
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        replacement_count = 0
        
        # Iterate through slides
        for slide in presentation.Slides:
            # Iterate through shapes
            for shape in slide.Shapes:
                # Check if shape has text
                if shape.HasTextFrame:
                    text_frame = shape.TextFrame
                    if text_frame.HasText:
                        for old_text, new_text in replacements.items():
                            if old_text in text_frame.TextRange.Text:
                                text_frame.TextRange.Text = text_frame.TextRange.Text.replace(old_text, new_text)
                                replacement_count += 1
                
                # Handle tables
                if shape.HasTable:
                    table = shape.Table
                    for row in range(1, table.Rows.Count + 1):
                        for col in range(1, table.Columns.Count + 1):
                            cell = table.Cell(row, col)
                            if cell.Shape.HasTextFrame and cell.Shape.TextFrame.HasText:
                                for old_text, new_text in replacements.items():
                                    if old_text in cell.Shape.TextFrame.TextRange.Text:
                                        cell.Shape.TextFrame.TextRange.Text = cell.Shape.TextFrame.TextRange.Text.replace(old_text, new_text)
                                        replacement_count += 1
        
        # Save and close
        presentation.SaveAs(output_path)
        presentation.Close()
        powerpoint.Quit()
        
        print(f"Successfully modified {input_file} -> {output_file}")
        print(f"Made {replacement_count} text replacements")
    
    except Exception as e:
        print(f"Error modifying .ppt file: {e}", file=sys.stderr)
        print("Note: .ppt files require PowerPoint to be installed on Windows", file=sys.stderr)
        sys.exit(1)


def main():
    parser = argparse.ArgumentParser(
        description="Modify text in PowerPoint presentations (.ppt and .pptx) based on a config file"
    )
    parser.add_argument(
        'input',
        help='Input PowerPoint file (.ppt or .pptx)'
    )
    parser.add_argument(
        '-o', '--output',
        help='Output file path (default: input_modified.ext)',
        default=None
    )
    parser.add_argument(
        '-c', '--config',
        help='Config file with text replacements (default: pptmodconfig.json)',
        default='pptmodconfig.json'
    )
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.exists(args.input):
        print(f"Error: Input file '{args.input}' not found", file=sys.stderr)
        sys.exit(1)
    
    # Determine output file
    if args.output is None:
        input_path = Path(args.input)
        args.output = str(input_path.parent / f"{input_path.stem}_modified{input_path.suffix}")
    
    # Load replacements from config
    replacements = load_config(args.config)
    
    if not replacements:
        print("Warning: No replacements found in config file", file=sys.stderr)
        sys.exit(1)
    
    print(f"Loaded {len(replacements)} replacement(s) from {args.config}")
    
    # Determine file type and process accordingly
    file_ext = Path(args.input).suffix.lower()
    
    if file_ext == '.pptx':
        modify_pptx(args.input, args.output, replacements)
    elif file_ext == '.ppt':
        modify_ppt(args.input, args.output, replacements)
    else:
        print(f"Error: Unsupported file type '{file_ext}'. Only .ppt and .pptx are supported", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
